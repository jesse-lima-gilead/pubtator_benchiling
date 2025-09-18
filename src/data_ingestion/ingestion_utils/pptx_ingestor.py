import datetime
from src.pubtator_utils.file_handler.base_handler import FileHandler
import xml.etree.ElementTree as ET
from pptx import Presentation
from xml.dom import minidom
from src.pubtator_utils.logs_handler.logger import SingletonLogger

# Initialize the logger
logger_instance = SingletonLogger()
logger = logger_instance.get_logger()


class PptxProcessor:
    """
    Processor for .pptx files to BioC XML.
    """

    def _prettify_xml(self, elem):
        """Return a pretty-printed XML string for the Element."""
        rough_string = ET.tostring(elem, "utf-8")
        reparsed = minidom.parseString(rough_string)
        return reparsed.toprettyxml(indent="  ")

    def _pptx_to_bioc(
        self,
        pptx_path,
        xml_path,
        file_handler: FileHandler,
        source_name="Gilead Internal Documents",
    ):
        # Load PPTX
        prs = Presentation(pptx_path)

        # Create root elements
        collection = ET.Element("collection")
        source = ET.SubElement(collection, "source")
        source.text = source_name
        date = ET.SubElement(collection, "date")
        date.text = datetime.date.today().isoformat()

        document = ET.SubElement(collection, "document")
        doc_id = ET.SubElement(document, "id")
        doc_id.text = pptx_path.split("/")[-1].split(".")[0]  # get the file name

        # Iterate slides
        for idx, slide in enumerate(prs.slides, start=1):
            # Get slide title or fallback
            title_shape = slide.shapes.title
            if title_shape and title_shape.text.strip():
                title_clean = title_shape.text.strip()
            else:
                title_clean = f"Slide {idx}"

            # Gather and clean all slide text (excluding title placeholder)
            texts = []
            for shape in slide.shapes:
                if shape is title_shape or not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in paragraph.runs)
                    if txt:
                        texts.append(txt)
            full_text = " ".join(texts)

            # Build passage
            passage = ET.SubElement(document, "passage")
            infon = ET.SubElement(passage, "infon", key="type")
            infon.text = title_clean
            offset = ET.SubElement(passage, "offset")
            offset.text = "0"
            text_el = ET.SubElement(passage, "text")
            text_el.text = full_text

        # Write pretty XML to file
        pretty_xml = self._prettify_xml(collection)
        file_handler.write_file(xml_path, pretty_xml)

    def _bioc_converter(
        self, input_doc_path: str, output_doc_dir: str, file_handler: FileHandler
    ):
        try:
            input_doc_name = input_doc_path.split("/")[-1]
            output_file_name = f'{input_doc_name.split(".")[0]}.xml'
            output_doc_path = f"{output_doc_dir}/{output_file_name}"
            self._pptx_to_bioc(input_doc_path, output_doc_path, file_handler)
        except Exception as ex:
            logger.exception(f"Failed to convert '{input_doc_path}' to BioC format")
            raise RuntimeError(f"Conversion error: {ex}") from ex

    def run(
        self,
        file_handler: FileHandler,
        internal_doc_name: str,
        internal_docs_path: str,
        bioc_path: str,
        **kwargs,
    ):
        try:
            input_doc_path = file_handler.get_file_path(
                internal_docs_path, internal_doc_name
            )
            self._bioc_converter(input_doc_path, bioc_path, file_handler)
            logger.info(
                f"Succesfully completed BioC file format conversion for doc: {internal_doc_name}"
            )
        except Exception as ex:
            logger.exception(
                f"Unexpected error during BioC run for doc: {internal_doc_name}"
            )
            raise RuntimeError(
                f"BioC run encountered an unexpected error: {ex}"
            ) from ex
