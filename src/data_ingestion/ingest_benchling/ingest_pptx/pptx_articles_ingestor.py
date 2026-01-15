"""
Benchling PPTX Articles Ingestor

Processes PPTX files from Benchling S3 bucket.
Converts PPTX -> BioC XML, extracts tables from slides.
Skips NER annotation steps.
"""

import os
from pathlib import Path
from typing import Any, Dict, Optional
import logging

from src.data_ingestion.ingest_benchling.benchling_config import BenchlingConfig
from src.pubtator_utils.file_handler.local_handler import LocalFileHandler

# Reuse existing PPTX processing functions
from src.data_ingestion.ingest_apollo.ingest_pptx.apollo_pptx_to_bioc_converter import (
    convert_pptx_to_bioc,
)
from src.data_ingestion.ingest_apollo.ingest_pptx.pptx_table_processor import (
    process_tables,
)

from pptx import Presentation
from pptx.exc import PackageNotFoundError
from datetime import datetime
import json

logger = logging.getLogger(__name__)


class BenchlingPPTXIngestor:
    """
    PPTX ingestor for Benchling documents.
    
    Pipeline:
    1. Read PPTX using python-pptx
    2. Extract text from slides
    3. Process tables from slides
    4. Convert to BioC XML
    5. Save metadata
    
    Note: Skips NER annotation as not applicable for Benchling.
    """
    
    def __init__(
        self,
        workflow_id: str,
        config: BenchlingConfig,
        paths: Dict[str, str],
        file_handler: LocalFileHandler,
        source: str = "benchling",
        write_to_delta: bool = True,
        **kwargs,
    ):
        self.workflow_id = workflow_id
        self.config = config
        self.paths = paths
        self.file_handler = file_handler
        self.source = source
        self.write_to_delta = write_to_delta
        
        self.ingestion_path = paths["ingestion_path"]
        self.interim_path = paths["ingestion_interim_path"]
        self.bioc_path = paths["bioc_path"]
        self.metadata_path = paths["metadata_path"]
        self.chunks_path = paths["chunks_path"]
        self.embeddings_path = paths["embeddings_path"]
        self.failed_path = paths["failed_ingestion_path"]
    
    def validate_pptx(self, pptx_path: str) -> bool:
        """Validate PPTX file can be opened."""
        try:
            Presentation(pptx_path)
            return True
        except (PackageNotFoundError, Exception) as e:
            logger.error(f"Invalid PPTX file {pptx_path}: {e}")
            return False
    
    def convert_to_bioc(self, file_name: str, metadata: Dict[str, Any]):
        """Convert PPTX to BioC XML."""
        try:
            convert_pptx_to_bioc(
                pptx_file_name=file_name,
                pptx_path=self.ingestion_path,
                bioc_path=self.bioc_path,
                metadata_path=self.metadata_path,
                pptx_interim_path=self.interim_path,
                article_metadata=metadata,
            )
            logger.info(f"Converted {file_name} to BioC XML")
            
        except Exception as e:
            logger.error(f"Failed to convert {file_name} to BioC: {e}")
            raise
    
    def extract_tables(
        self,
        pptx_path: str,
        document_id: str,
        metadata: Dict[str, Any],
    ) -> list:
        """Extract tables from PPTX file."""
        try:
            prs = Presentation(pptx_path)
            all_table_details = []
            
            for slide_idx, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        
                        # Extract table data
                        table_data = []
                        for row in table.rows:
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text)
                            table_data.append(row_data)
                        
                        # Create flat text representation
                        flat_text = "\n".join([
                            " | ".join(row) for row in table_data
                        ])
                        
                        table_detail = {
                            "payload": {
                                "document_grsar_id": document_id,
                                "slide_idx": slide_idx,
                                "table_id": f"{document_id}_slide{slide_idx}_table",
                                "row_count": len(table.rows),
                                "column_count": len(table.columns),
                                "clean_flat_text": flat_text,
                                "context_flat_text": flat_text,
                                **metadata,
                            }
                        }
                        all_table_details.append(table_detail)
            
            logger.info(f"Extracted {len(all_table_details)} tables from {document_id}")
            return all_table_details
            
        except Exception as e:
            logger.error(f"Failed to extract tables from {pptx_path}: {e}")
            return []
    
    def extract_metadata(self, file_name: str, file_path: str) -> Dict[str, Any]:
        """Extract metadata from PPTX file."""
        document_grsar_id = file_name.split(".")[0]
        
        metadata = {
            "document_grsar_id": document_grsar_id,
            "source": self.source,
            "file_name": file_name,
            "file_path": file_path,
            "workflow_id": self.workflow_id,
            "processing_date": datetime.now().isoformat(),
            "file_type": "pptx",
        }
        
        # Try to get file size and slide count
        full_path = os.path.join(self.ingestion_path, file_name)
        if os.path.exists(full_path):
            metadata["size_bytes"] = os.path.getsize(full_path)
            try:
                prs = Presentation(full_path)
                metadata["slide_count"] = len(prs.slides)
            except:
                pass
        
        return metadata
    
    def save_metadata(self, metadata: Dict[str, Any], document_id: str):
        """Save metadata to JSON file."""
        Path(self.metadata_path).mkdir(parents=True, exist_ok=True)
        metadata_file = os.path.join(self.metadata_path, f"{document_id}_metadata.json")
        
        with open(metadata_file, "w", encoding="utf-8") as f:
            json.dump(metadata, f, indent=2, ensure_ascii=False)
        
        logger.info(f"Saved metadata to {metadata_file}")
    
    def run(self, file_name: str):
        """
        Run the PPTX ingestion pipeline.
        
        Args:
            file_name: Name of PPTX file in ingestion directory
        """
        if not file_name.lower().endswith((".pptx", ".ppt")):
            logger.warning(f"{file_name} is not a PPTX file")
            return
        
        logger.info(f"Processing PPTX: {file_name}")
        
        document_id = file_name.split(".")[0]
        pptx_path = os.path.join(self.ingestion_path, file_name)
        
        try:
            # Validate file
            if not self.validate_pptx(pptx_path):
                raise Exception(f"Invalid PPTX file: {file_name}")
            
            # Step 1: Extract metadata
            metadata = self.extract_metadata(file_name, pptx_path)
            self.save_metadata(metadata, document_id)
            
            # Step 2: Extract tables
            logger.info(f"Extracting tables from {file_name}")
            table_details = self.extract_tables(pptx_path, document_id, metadata)
            
            # Save table details
            if table_details:
                Path(self.embeddings_path).mkdir(parents=True, exist_ok=True)
                tables_file = os.path.join(self.embeddings_path, f"{document_id}_tables.json")
                with open(tables_file, "w", encoding="utf-8") as f:
                    json.dump(table_details, f, indent=2, ensure_ascii=False)
            
            # Step 3: Convert to BioC XML
            logger.info(f"Converting to BioC XML")
            self.convert_to_bioc(file_name, metadata)
            
            logger.info(f"Successfully processed PPTX: {file_name}")
            
        except Exception as e:
            logger.error(f"Failed to process PPTX {file_name}: {e}")
            
            # Move to failed directory
            failed_path = os.path.join(self.failed_path, file_name)
            if os.path.exists(pptx_path):
                Path(self.failed_path).mkdir(parents=True, exist_ok=True)
                self.file_handler.move_file(pptx_path, failed_path)
            
            raise
