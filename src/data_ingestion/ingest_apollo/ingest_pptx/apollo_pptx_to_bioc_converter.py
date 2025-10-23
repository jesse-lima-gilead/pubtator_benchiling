import datetime
import json
import re
from pathlib import Path
from typing import Optional

from pptx import Presentation
import bioc
from src.pubtator_utils.file_handler.base_handler import FileHandler
from src.pubtator_utils.logs_handler.logger import SingletonLogger

# Initialize logger
logger_instance = SingletonLogger()
logger = logger_instance.get_logger()


class PptxProcessor:
    """
    PPTX -> BioC converter.

    Key guarantees (preserve original behaviour):
      - Do NOT include table text
      - Join all slide text into a single line (space-separated)
      - passage.offset is always 0
      - Only first-level shapes are considered (no recursive group traversal)
      - Source name default: "Gilead Internal: Apollo - pptx"
      - Optional metadata_fields are written into document.infons
      - Collection date is current date YYYY-MM-DD
    """

    def clean_xml_text(self, s):
        _RE_XML_ILLEGAL = re.compile(r"[\x00-\x08\x0B\x0C\x0E-\x1F\x7F-\x9F]")
        if s is None:
            return ""
        if not isinstance(s, str):
            s = str(s)
        return _RE_XML_ILLEGAL.sub("", s)

    def _collect_text_from_shape(self, shape):
        """Collect trimmed text blocks from a first-level shape.
        Tables are skipped explicitly. Group shapes are not traversed.
        """
        blocks = []

        # shapes with text frame
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            try:
                for paragraph in shape.text_frame.paragraphs:
                    run_text = "".join(
                        (r.text or "") for r in getattr(paragraph, "runs", [])
                    )
                    if not run_text:
                        run_text = getattr(paragraph, "text", "") or ""
                    if run_text and run_text.strip():
                        blocks.append(run_text.strip())
            except Exception:
                s = getattr(shape, "text", "")
                if s and s.strip():
                    blocks.append(s.strip())

        # explicitly skip tables
        elif getattr(shape, "has_table", False) and shape.has_table:
            logger.debug("Skipping table text as per config.")

        else:
            # fallback for other shapes (including groups) - only top-level text
            s = getattr(shape, "text", "")
            if s and s.strip():
                blocks.append(s.strip())

        return blocks

    def _convert_pptx_to_bioc(
        self,
        pptx_path: str,
        xml_path: str,
        file_handler: FileHandler,
        write_to_s3: bool,
        s3_bioc_path: str,
        s3_file_handler: FileHandler,
        source_name: str = "Gilead Internal: Apollo - pptx",
        metadata_fields: dict = None,
        include_notes: bool = False,
    ):
        """Core conversion logic. Raises on errors.
        Behaviour mirrors the original implementation but consolidated into one method.
        """
        prs = Presentation(pptx_path)

        coll = bioc.BioCCollection()
        coll.source = source_name
        coll.date = datetime.date.today().isoformat()

        doc = bioc.BioCDocument()
        doc.id = Path(pptx_path).stem

        # attach metadata fields
        if metadata_fields and isinstance(metadata_fields, dict):
            for k, v in metadata_fields.items():
                doc.infons[k] = "" if v is None else str(v)

        for idx, slide in enumerate(prs.slides, start=1):
            # collect first-level shapes
            blocks = []
            for shape in slide.shapes:
                blocks.extend(self._collect_text_from_shape(shape))

            # append notes if requested
            if (
                include_notes
                and getattr(slide, "has_notes_slide", False)
                and slide.has_notes_slide
            ):
                try:
                    notes_tf = slide.notes_slide.notes_text_frame
                    notes_text = notes_tf.text if notes_tf is not None else ""
                    if notes_text and notes_text.strip():
                        blocks.append(notes_text.strip())
                except Exception:
                    logger.debug(
                        f"Unable to read notes for slide {idx}; skipping notes."
                    )

            # normalize and join into single-line passage
            normalized_blocks = [" ".join(b.split()) for b in blocks if b.strip()]
            full_text = " ".join(normalized_blocks).strip()

            passage = bioc.BioCPassage()
            passage.offset = 0
            passage.infons["type"] = "section"
            passage.infons["section_title"] = str(idx)
            passage.text = self.clean_xml_text(full_text)
            doc.passages.append(passage)

        coll.documents.append(doc)

        # merge small passages
        col_with_merged_passages = merge_small_passages_in_collection(
            collection=coll,
            threshold_words=75,
            max_iterations=5,
            prefer_merge_with_next=True,
        )

        # write local XML via provided file_handler
        file_handler.write_file_as_bioc(xml_path, col_with_merged_passages)
        logger.info(f"Wrote BioC XML for {pptx_path} -> {xml_path}")

        # # optionally write to S3
        # if write_to_s3:
        #     bioc_file_name = Path(xml_path).name
        #     s3_file_path = s3_file_handler.get_file_path(s3_bioc_path, bioc_file_name)
        #     s3_file_handler.write_file_as_bioc(s3_file_path, col_with_merged_passages)
        #     logger.info(f"Saving BioC XML to S3: {s3_file_path}")

    def run(
        self,
        file_handler: FileHandler,
        internal_doc_name: str,
        internal_docs_path: str,
        bioc_path: str,
        write_to_s3: bool,
        s3_bioc_path: str,
        s3_file_handler: FileHandler,
        metadata_fields: dict = None,
        **kwargs,
    ):
        """Public runner that prepares paths and invokes conversion.
        Keeps behaviour and signature compatible with the original code.
        """
        try:
            input_doc_path = file_handler.get_file_path(
                internal_docs_path, internal_doc_name
            )
            input_doc_name = Path(input_doc_path).name
            output_file_name = f"{Path(input_doc_name).stem}.xml"
            output_doc_path = file_handler.get_file_path(bioc_path, output_file_name)

            include_notes = kwargs.get("include_notes", False)

            self._convert_pptx_to_bioc(
                pptx_path=input_doc_path,
                xml_path=output_doc_path,
                file_handler=file_handler,
                write_to_s3=write_to_s3,
                s3_bioc_path=s3_bioc_path,
                s3_file_handler=s3_file_handler,
                metadata_fields=metadata_fields,
                include_notes=include_notes,
            )

            logger.info(
                f"Successfully completed BioC file format conversion for doc: {internal_doc_name}"
            )
        except Exception as ex:
            logger.exception(
                f"Unexpected error during BioC run for doc: {internal_doc_name}"
            )
            raise RuntimeError(
                f"BioC run encountered an unexpected error: {ex}"
            ) from ex


def pptx_to_bioc_converter(
    file_handler: FileHandler,
    internal_doc_name: str,
    internal_docs_path: str,
    bioc_path: str,
    metadata_fields: dict,
    write_to_s3: bool,
    s3_bioc_path: str,
    s3_file_handler: FileHandler,
):
    """Simple module-level wrapper kept for compatibility with existing callers."""
    processor = PptxProcessor()
    processor.run(
        file_handler=file_handler,
        internal_doc_name=internal_doc_name,
        internal_docs_path=internal_docs_path,
        bioc_path=bioc_path,
        write_to_s3=write_to_s3,
        s3_bioc_path=s3_bioc_path,
        s3_file_handler=s3_file_handler,
        metadata_fields=metadata_fields,
    )


def _normalize_whitespace_single_line(s: Optional[str]) -> str:
    """Normalize whitespace to a single line (no newlines), collapse multiple spaces to one."""
    if not s:
        return ""
    # Replace NBSP and other non-breaking spaces with normal space
    s = s.replace("\u00A0", " ")
    s = s.replace("\u200B", "")  # zero-width space -> remove
    # convert \r\n,\r,\n to spaces
    s = s.replace("\r\n", " ").replace("\r", " ").replace("\n", " ")
    # collapse whitespace
    s = re.sub(r"[ \t\r\n]+", " ", s)
    return s.strip()


def _try_merge_provenance(a: Optional[str], b: Optional[str]) -> Optional[str]:
    """Try to sensibly merge two provenance infons (JSON-aware), return JSON string or joined string."""
    if not a and not b:
        return None
    if a and not b:
        return a
    if b and not a:
        return b
    try:
        va = json.loads(a)
    except Exception:
        va = a
    try:
        vb = json.loads(b)
    except Exception:
        vb = b

    if isinstance(va, list) and isinstance(vb, list):
        combined = va + vb
    elif isinstance(va, list):
        combined = va + [vb]
    elif isinstance(vb, list):
        combined = [va] + vb
    else:
        combined = [va, vb]

    try:
        return json.dumps(combined, ensure_ascii=False)
    except Exception:
        return " | ".join(
            [
                _normalize_whitespace_single_line(str(va)),
                _normalize_whitespace_single_line(str(vb)),
            ]
        )


_EXEC_SUMMARY_RE = re.compile(r"^\s*executive\s+summary\s*$", re.IGNORECASE)


def _is_executive_summary_title(title: str) -> bool:
    """Regex-based check for 'executive summary' section titles."""
    if not title:
        return False
    return bool(_EXEC_SUMMARY_RE.match(title.strip()))


def _get_section_title_from_passage(passage) -> str:
    return passage.infons.get("section_title", "") if passage.infons else ""


def merge_small_passages_in_collection(
    collection: bioc.BioCCollection,
    threshold_words: int = 100,
    max_iterations: int = 5,
    prefer_merge_with_next: bool = True,
) -> bioc.BioCCollection:
    for doc in collection.documents:
        iteration = 0
        changed = True

        while changed and iteration < max_iterations:
            iteration += 1
            changed = False
            old_passages = list(doc.passages)
            n = len(old_passages)
            i = 0
            new_passages = []

            while i < n:
                cur = old_passages[i]
                cur_text = _normalize_whitespace_single_line(cur.text)
                cur_count = len(cur_text.split())

                if cur_count >= threshold_words:
                    p = bioc.BioCPassage()
                    p.offset = getattr(cur, "offset", 0)
                    p.text = cur_text
                    p.infons = dict(cur.infons) if cur.infons else {}
                    p.annotations = (
                        list(cur.annotations) if hasattr(cur, "annotations") else []
                    )
                    p.relations = (
                        list(cur.relations) if hasattr(cur, "relations") else []
                    )
                    new_passages.append(p)
                    i += 1
                    continue

                # Start a run of small passages
                run_texts = [cur_text]
                run_indices = [i]
                run_titles = [cur.infons.get("section_title", "") if cur.infons else ""]
                j = i + 1
                while j < n:
                    nxt = old_passages[j]
                    nxt_text = _normalize_whitespace_single_line(nxt.text)
                    nxt_count = len(nxt_text.split())
                    if nxt_count < threshold_words:
                        run_texts.append(nxt_text)
                        run_indices.append(j)
                        run_titles.append(
                            nxt.infons.get("section_title", "") if nxt.infons else ""
                        )
                        j += 1
                        continue
                    break
                run_merged_text = _normalize_whitespace_single_line(
                    " ".join([t for t in run_texts if t])
                )
                run_word_count = len(run_merged_text.split())
                merged_section_title = (
                    " | ".join([t for t in run_titles if t]).strip() or "body_content"
                )

                if run_word_count >= threshold_words:
                    base = old_passages[run_indices[-1]]
                    p = bioc.BioCPassage()
                    p.offset = getattr(base, "offset", 0)
                    p.text = run_merged_text
                    p.infons = dict(base.infons) if base.infons else {}
                    p.infons["section_title"] = merged_section_title
                    p.infons["type"] = p.infons.get("type", "section")
                    prov = None
                    for idx in run_indices:
                        prov = _try_merge_provenance(
                            prov,
                            old_passages[idx].infons.get("provenance")
                            if old_passages[idx].infons
                            else None,
                        )
                    if prov:
                        p.infons["provenance"] = prov
                    p.annotations = []
                    p.relations = []
                    for idx in run_indices:
                        if (
                            hasattr(old_passages[idx], "annotations")
                            and old_passages[idx].annotations
                        ):
                            p.annotations.extend(old_passages[idx].annotations)
                        if (
                            hasattr(old_passages[idx], "relations")
                            and old_passages[idx].relations
                        ):
                            p.relations.extend(old_passages[idx].relations)
                    new_passages.append(p)
                    i = j
                    changed = True
                    continue

                next_idx = j if j < n else None
                prev_exists = len(new_passages) > 0

                if (
                    prefer_merge_with_next
                    and next_idx is not None
                    and not _is_executive_summary_title(
                        _get_section_title_from_passage(old_passages[next_idx])
                    )
                ):
                    nxt = old_passages[next_idx]
                    nxt_text = _normalize_whitespace_single_line(nxt.text)
                    merged_text = _normalize_whitespace_single_line(
                        run_merged_text + " " + nxt_text
                    )
                    base = nxt
                    p = bioc.BioCPassage()
                    p.offset = getattr(base, "offset", 0)
                    p.text = merged_text
                    p.infons = dict(base.infons) if base.infons else {}
                    # Collect all section titles from run and next
                    all_titles = [t for t in run_titles if t]
                    next_title = _get_section_title_from_passage(nxt)
                    if next_title:
                        all_titles.append(next_title)
                    merged_all_titles = " | ".join(all_titles).strip() or "body_content"
                    p.infons["section_title"] = merged_all_titles
                    p.infons["type"] = p.infons.get("type", "section")
                    prov = None
                    for idx in run_indices:
                        prov = _try_merge_provenance(
                            prov,
                            old_passages[idx].infons.get("provenance")
                            if old_passages[idx].infons
                            else None,
                        )
                    prov = _try_merge_provenance(
                        prov, nxt.infons.get("provenance") if nxt.infons else None
                    )
                    if prov:
                        p.infons["provenance"] = prov
                    p.annotations = []
                    p.relations = []
                    for idx in run_indices:
                        if (
                            hasattr(old_passages[idx], "annotations")
                            and old_passages[idx].annotations
                        ):
                            p.annotations.extend(old_passages[idx].annotations)
                        if (
                            hasattr(old_passages[idx], "relations")
                            and old_passages[idx].relations
                        ):
                            p.relations.extend(old_passages[idx].relations)
                    if hasattr(nxt, "annotations") and nxt.annotations:
                        p.annotations.extend(nxt.annotations)
                    if hasattr(nxt, "relations") and nxt.relations:
                        p.relations.extend(nxt.relations)

                    new_passages.append(p)
                    i = next_idx + 1
                    changed = True
                    continue

                if prev_exists:
                    prev = new_passages.pop()
                    prev_text = _normalize_whitespace_single_line(prev.text)
                    merged_text = _normalize_whitespace_single_line(
                        prev_text + " " + run_merged_text
                    )
                    merged_infons = dict(prev.infons) if prev.infons else {}
                    # Collect all section titles from prev and run
                    prev_title = (
                        prev.infons.get("section_title", "") if prev.infons else ""
                    )
                    all_titles = [prev_title] + [t for t in run_titles if t]
                    merged_all_titles = (
                        " | ".join([t for t in all_titles if t]).strip()
                        or "body_content"
                    )
                    merged_infons["section_title"] = merged_all_titles
                    merged_infons["type"] = merged_infons.get("type", "section")
                    prov = _try_merge_provenance(
                        prev.infons.get("provenance") if prev.infons else None, None
                    )
                    for idx in run_indices:
                        prov = _try_merge_provenance(
                            prov,
                            old_passages[idx].infons.get("provenance")
                            if old_passages[idx].infons
                            else None,
                        )
                    if prov:
                        merged_infons["provenance"] = prov
                    merged_annotations = (
                        list(prev.annotations) if hasattr(prev, "annotations") else []
                    )
                    merged_relations = (
                        list(prev.relations) if hasattr(prev, "relations") else []
                    )
                    for idx in run_indices:
                        if (
                            hasattr(old_passages[idx], "annotations")
                            and old_passages[idx].annotations
                        ):
                            merged_annotations.extend(old_passages[idx].annotations)
                        if (
                            hasattr(old_passages[idx], "relations")
                            and old_passages[idx].relations
                        ):
                            merged_relations.extend(old_passages[idx].relations)
                    p = bioc.BioCPassage()
                    p.offset = getattr(prev, "offset", 0)
                    p.text = merged_text
                    p.infons = merged_infons
                    p.annotations = merged_annotations
                    p.relations = merged_relations
                    new_passages.append(p)
                    i = j
                    changed = True
                    continue

                p = bioc.BioCPassage()
                p.offset = getattr(old_passages[run_indices[-1]], "offset", 0)
                p.text = run_merged_text
                p.infons = (
                    dict(old_passages[run_indices[-1]].infons)
                    if old_passages[run_indices[-1]].infons
                    else {}
                )
                p.infons["section_title"] = merged_section_title
                p.infons["type"] = p.infons.get("type", "section")
                p.annotations = []
                p.relations = []
                for idx in run_indices:
                    if (
                        hasattr(old_passages[idx], "annotations")
                        and old_passages[idx].annotations
                    ):
                        p.annotations.extend(old_passages[idx].annotations)
                    if (
                        hasattr(old_passages[idx], "relations")
                        and old_passages[idx].relations
                    ):
                        p.relations.extend(old_passages[idx].relations)
                new_passages.append(p)
                i = j

            # End of while i < n
            doc.passages = new_passages

    return collection
