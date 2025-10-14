import re
import uuid
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd

from src.data_ingestion.ingestion_utils.s3_uploader import upload_to_s3
from src.pubtator_utils.file_handler.base_handler import FileHandler
from src.pubtator_utils.logs_handler.logger import SingletonLogger

# Initialize logger
logger_instance = SingletonLogger()
logger = logger_instance.get_logger()


class PptxTableExtractor:
    def __init__(
        self,
        pptx_path: str,
        interim_excel_path: str,
        embeddings_output_dir: str,
        file_handler: FileHandler,
        write_to_s3: bool,
        s3_embeddings_path: str,
        s3_interim_path: str,
        s3_file_handler: FileHandler,
    ):
        self.pptx_path = pptx_path
        self.source_filename = Path(pptx_path).stem
        self.interim_path = Path(interim_excel_path)
        self.embeddings_output_dir = Path(embeddings_output_dir)
        self.file_handler = file_handler
        self.write_to_s3 = write_to_s3
        self.s3_embeddings_path = s3_embeddings_path
        self.s3_interim_path = s3_interim_path
        self.s3_file_handler = s3_file_handler
        self.interim_path.mkdir(exist_ok=True, parents=True)
        self.embeddings_output_dir.mkdir(exist_ok=True, parents=True)
        self.prs = Presentation(pptx_path)
        self.tables_data = []

    def _collect_text_from_shape(self, shape):
        """Extract trimmed text blocks from first-level shapes, skip tables explicitly."""
        blocks = []
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            try:
                for paragraph in shape.text_frame.paragraphs:
                    run_text = "".join(
                        (r.text or "") for r in getattr(paragraph, "runs", [])
                    )
                    if not run_text:
                        run_text = getattr(paragraph, "text", "") or ""
                    if run_text and run_text.strip():
                        blocks.append(run_text.strip())
            except Exception:
                s = getattr(shape, "text", "")
                if s and s.strip():
                    blocks.append(s.strip())
        elif getattr(shape, "has_table", False) and shape.has_table:
            pass  # skip table text
        else:
            s = getattr(shape, "text", "")
            if s and s.strip():
                blocks.append(s.strip())
        return blocks

    def _collect_slide_text(self, slide):
        blocks = []
        for shape in slide.shapes:
            blocks.extend(self._collect_text_from_shape(shape))
        normalized_blocks = [" ".join(b.split()) for b in blocks if b.strip()]
        return " ".join(normalized_blocks).strip()

    def _extract_table_id_and_name_from_slide(self, slide_text: str, idx: int):
        table_id = f"Extracted_Table_{idx}"
        table_name = f"Extracted_Table_{idx}"
        candidate = (slide_text or "").strip()
        if candidate:
            match = re.search(
                r"\b(?:Table|TABLE|Tbl|tbl)[\s\.\-]*([0-9]+[a-zA-Z]?)\b", candidate
            )
            if match:
                num = match.group(1)
                table_id = f"Table_{num}"
                table_name = candidate.strip()
            else:
                if len(candidate.split()) <= 15:
                    table_name = candidate.strip()

        table_name = re.sub(r"[^A-Za-z0-9]", "_", table_name)
        # Maximum length for table_name in the filename
        max_table_name_len = 25

        # Truncate table_name if it's too long
        if len(table_name) > max_table_name_len:
            table_name = table_name[:max_table_name_len]

        return table_id, table_name

    def _process_table(
        self, shape, slide_idx, tbl_idx, slide_text, metadata_fields: dict = None
    ):
        table = shape.table
        data = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        if not data:
            return None

        df = pd.DataFrame(data[1:], columns=data[0])

        header_text = " ".join([str(h) for h in df.columns if h])
        data_text = " ".join(
            [str(cell) for row in df.values.tolist() for cell in row if cell]
        )
        flattened_text = " ".join([header_text, data_text])

        table_id, table_name = self._extract_table_id_and_name_from_slide(
            slide_text, tbl_idx
        )
        table_html = df.to_html(index=False, escape=False)

        # Excel path: interim_path / article_id / table_name / file.xlsx
        excel_dir = self.interim_path / self.source_filename
        excel_dir.mkdir(parents=True, exist_ok=True)
        excel_file = (
            excel_dir / f"{self.source_filename}_slide_{slide_idx}_{table_name}.xlsx"
        )
        df.to_excel(excel_file, index=False)
        logger.info(f"Table extarcted: {excel_file}")

        table_dict = {
            "table_sequence": tbl_idx,
            "slide_index": slide_idx,
            "table_id": str(uuid.uuid4()),
            "chunk_processing_date": datetime.now().date().isoformat(),
            "article_id": self.source_filename,
            "article_table_id": table_id,
            "table_name": table_name,
            "slide_text": slide_text,
            "chunk_type": "table_chunk",
            "processing_ts": datetime.now().isoformat(),
            "columns": list(df.columns.astype(str)),
            "row_count": df.shape[0],
            "column_count": df.shape[1],
            "clean_flat_text": slide_text + " " + flattened_text,
            "merged_text": slide_text + " " + table_html,
            "excel_path": str(excel_file),
        }

        # Merge optional metadata_fields
        if metadata_fields and isinstance(metadata_fields, dict):
            table_dict.update(metadata_fields)

        return table_dict

    def extract_tables(self, metadata_fields: dict = None):
        tbl_counter = 1
        for slide_idx, slide in enumerate(self.prs.slides, start=1):
            slide_text = self._collect_slide_text(slide)
            for shape in slide.shapes:
                # Slide-level tables
                if getattr(shape, "has_table", False) and shape.has_table:
                    tdata = self._process_table(
                        shape, slide_idx, tbl_counter, slide_text, metadata_fields
                    )
                    if tdata:
                        self.tables_data.append(tdata)
                        tbl_counter += 1
                # Tables inside group shapes (1-level)
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for child in shape.shapes:
                        if getattr(child, "has_table", False) and child.has_table:
                            tdata = self._process_table(
                                child,
                                slide_idx,
                                tbl_counter,
                                slide_text,
                                metadata_fields,
                            )
                            if tdata:
                                self.tables_data.append(tdata)
                                tbl_counter += 1

        if self.write_to_s3 and tbl_counter > 1:
            pptx_interim_file_upload_counter = 0
            ingestion_interim_path = self.file_handler.get_file_path(
                self.interim_path, self.source_filename
            )
            s3_ingestion_interim_path = self.s3_file_handler.get_file_path(
                self.s3_interim_path, self.source_filename
            )
            for pptx_interim_file in self.file_handler.list_files(
                ingestion_interim_path
            ):
                local_file_path = self.file_handler.get_file_path(
                    ingestion_interim_path, pptx_interim_file
                )
                s3_file_path = self.s3_file_handler.get_file_path(
                    s3_ingestion_interim_path, pptx_interim_file
                )
                logger.info(
                    f"Uploading file {pptx_interim_file} to S3 path {s3_file_path}"
                )
                upload_to_s3(
                    local_path=local_file_path,
                    s3_path=s3_file_path,
                    s3_file_handler=self.s3_file_handler,
                )
                pptx_interim_file_upload_counter += 1
            logger.info(
                f"For {self.source_filename}.pptx, Total Interim Files uploaded to S3: {pptx_interim_file_upload_counter}"
            )

        return self.tables_data

    def write_table_embeddings(self):
        if not self.tables_data:
            return None
        json_list = [{"payload": t} for t in self.tables_data]
        file_name = f"{self.source_filename}_tables.json"
        json_file = self.file_handler.get_file_path(
            self.embeddings_output_dir, file_name
        )
        self.file_handler.write_file_as_json(json_file, json_list)
        logger.info(f"Written {self.source_filename} table embeddings to {json_file}")

        if self.write_to_s3:
            s3_file_path = self.s3_file_handler.get_file_path(
                self.s3_embeddings_path, file_name
            )
            self.s3_file_handler.write_file_as_json(s3_file_path, json_list)
            logger.info(f"Saving table embeddings to S3: {s3_file_path}")


def extract_pptx_tables(
    file_handler: FileHandler,
    pptx_path: str,
    interim_dir: str,
    embeddings_dir: str,
    bioc_metadata_fields: dict,
    write_to_s3: bool,
    s3_embeddings_path: str,
    s3_interim_path: str,
    s3_file_handler: FileHandler,
):
    """
    Extracts tables from a PPTX file, writes them to a JSON Embeddings payload.
    """
    extractor = PptxTableExtractor(
        pptx_path,
        interim_dir,
        embeddings_dir,
        file_handler,
        write_to_s3,
        s3_embeddings_path,
        s3_interim_path,
        s3_file_handler,
    )
    tables = extractor.extract_tables(bioc_metadata_fields)
    logger.info(f"Extracted {len(tables)} tables from {pptx_path}")

    extractor.write_table_embeddings()

    logger.info(f"Table Extraction completed for {pptx_path}")
