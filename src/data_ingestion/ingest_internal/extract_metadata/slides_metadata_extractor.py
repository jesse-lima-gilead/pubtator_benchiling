from pptx import Presentation
from pathlib import Path
from typing import Union, Dict, Any


class SlidesMetadataExtractor:
    def extract_metadata(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        prs = Presentation(file_path)
        props = prs.core_properties
        meta: Dict[str, Any] = {
            "title": props.title,
            "author": props.author,
            "created": props.created,
            "modified": props.modified,
            "subject": props.subject,
            "keywords": props.keywords,
            "slides": [],
        }
        # Slide‑level titles and notes as “abstract”
        for slide in prs.slides:
            title = slide.shapes.title.text if slide.shapes.title else None
            notes = (
                slide.notes_slide.notes_text_frame.text
                if slide.has_notes_slide
                else None
            )
            meta["slides"].append({"title": title, "notes": notes})
        return meta
